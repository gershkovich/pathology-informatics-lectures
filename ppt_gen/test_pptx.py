#!/usr/bin/env python3

print("Testing python-pptx")
print("==================")

try:
    from pptx import Presentation
    print("✓ Successfully imported python-pptx")
    
    # Create a simple presentation
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Created with python-pptx"
    prs.save("test.pptx")
    print("✓ Successfully created test.pptx")
except ImportError as e:
    print("✗ Error importing python-pptx:", e)
except Exception as e:
    print("✗ Error creating presentation:", e)

print("\nIf test.pptx was created, your setup is working correctly.")
