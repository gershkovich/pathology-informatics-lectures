from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Define theme colors
COLOR_PRIMARY = RGBColor(0, 51, 102)  # Deep Blue
COLOR_ACCENT = RGBColor(0, 184, 217)  # Light Blue

# Font sizes
TITLE_FONT_SIZE = Pt(44)
SUBTITLE_FONT_SIZE = Pt(24)
BODY_FONT_SIZE = Pt(20)

# Helper to add slides
def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_content_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = BODY_FONT_SIZE

# Slide Content
add_title_slide(prs, "AI-Powered Dialogue System for Diagnostic Quality Management",
                "Reducing Errors | Enhancing Efficiency | Transforming Pathology Workflow")

add_content_slide(prs, "The Problem – Inefficiency & Risk in Diagnostics", [
    "Sentinel Events caused by misaligned reports and missed discrepancies",
    "Amended pathology reports delay care, increase cost, and erode trust",
    "Manual document review is time-consuming and error-prone",
    "Triage errors, improper scheduling, and inefficient care pathways persist"
])

add_content_slide(prs, "Our Solution – Dialogue Project", [
    "Real-time AI-based Quality Control for documentation alignment",
    "Detects errors and misalignments across pathology, radiology, clinical notes",
    "Prevents incorrect treatments and improper scheduling",
    "Minimizes the need for report amendments and streamlines reporting"
])

add_content_slide(prs, "Key Impact Areas", [
    "Prevent Sentinel Events by highlighting missed critical findings",
    "Improve Pathology Quality by reducing report amendments",
    "Enhance Radiology Alignment with pathology and clinical data",
    "Enable Efficient Triage and resource allocation",
    "Reduce Waste and administrative overhead in healthcare delivery"
])

add_content_slide(prs, "Technology Overview", [
    "Web-based AI Workspace built for Pathology Informatics",
    "Java 17, Spring Boot 3.x, integrated with Yale AWS infrastructure",
    "AI algorithms for report alignment, error detection, and workflow optimization",
    "Scalable, secure, and adaptable across healthcare environments"
])

add_content_slide(prs, "Transforming Diagnostic Quality", [
    "Revolutionizing Diagnostic Quality with AI-driven tools",
    "Integrated within Yale’s innovation framework",
    "DQMS (Diagnostic Quality Management System) to reduce errors at scale",
    "Fewer misdiagnoses = better patient outcomes and safer healthcare",
    "Positioned to set a new standard for pathology quality assurance"
])

add_content_slide(prs, "Traction & Next Steps", [
    "Prototype Ready: Real-time error detection and workflow enhancement",
    "Next Milestones: Clinical environment hardening, pilot with teams",
    "Secure funding for scale-up & market launch",
    "Seeking Support: Resources, partnerships, and strategic input"
])

# Save locally
prs.save("Dialogue_Project_Pitch_Deck.pptx")