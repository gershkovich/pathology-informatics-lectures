from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
slide_width, slide_height = prs.slide_width, prs.slide_height

# Helper to add title and content
def add_slide(title, content_lines, bullet=True):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    tf = body_placeholder.text_frame
    tf.clear()
    for line in content_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        if bullet:
            p.level = 0
        else:
            p.level = 1

# Slide 1
add_slide(
    "Why Now? Clinical-grade AI demands clinical-grade software",
    [
        "Independent AI decisions now inevitable (FDA, EU AI Act)",
        "Traditional dev cycles too slow for model change",
        "Regulated healthcare demands audit, security, traceability",
        "\nEvery software defect becomes a potential patient-safety event"
    ]
)

# Slide 2
add_slide(
    "The Missing Piece: AI-Assisted SDLC",
    [
        "LLMs are not just engines—they're pit crews.",
        "Pipeline: Code Commit → LLM Code Gen + Patch → Static/Dynamic Analysis → Human + LLM Review → Signed Artifact → Deploy → Monitor",
        "- Auto-generate tests\n- Inline risk annotation\n- SBOM and dependency scan"
    ]
)

# Slide 3
add_slide(
    "Speed + Quality: Two Case Studies",
    [
        "Project A (2024) vs Project B (2025):",
        "Tooling: Cloud+ChatGPT vs Windsurf+IntelliJ AI",
        "Time to MVP: 4w vs 6d",
        "LOC:Test Ratio: 1:0.4 vs 1:1.1",
        "Bugs found: 8 vs 31 (AI 18 + human 13)",
        "Regulatory: Trace Matrix vs Trace+SBOM+Threat Model"
    ]
)

# Slide 4
add_slide(
    "How Hedgehog Catches What Humans Miss",
    [
        "Hybrid Review = Amplified Vigilance",
        "Security: Logic edge-case caught by GPT-4",
        "Safety: Dependency CVE auto-flagged",
        "Maintainability: Unsafe regex found, corrected by LLM",
        "Performance: -42% latency after LLM optimization"
    ]
)

# Slide 5
add_slide(
    "Roadmap: From CI/CD to CQMS",
    [
        "Today: Automated tests, SCA, basic risk logging",
        "Next: Live risk scoring, validation auto-gen",
        "Future: Drift monitoring, self-updating IEC 62304 docs, AI-prepared FDA drafts",
        "\n\"14 days → 3 hours median change-to-production. 0 Sev-1 bugs in 12 months.\""
    ]
)

# Save presentation
prs.save('../data/hedgehog_ai_sdcl_poster.pptx')
