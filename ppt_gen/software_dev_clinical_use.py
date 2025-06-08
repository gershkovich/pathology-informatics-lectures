from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Define theme colors
COLOR_PRIMARY = RGBColor(0, 51, 102)  # Deep Blue
COLOR_ACCENT = RGBColor(0, 184, 217)  # Light Blue

# Font sizes
TITLE_FONT_SIZE = Pt(44)
SUBTITLE_FONT_SIZE = Pt(24)
BODY_FONT_SIZE = Pt(20)

# Helper to add slides
def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_content_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = BODY_FONT_SIZE

# Slide Content based on custom_development.md
add_title_slide(prs, "Software Development for Clinical Use in Pathology",
                "Peter Gershkovich, MD. MHA\nYale University School of Medicine | June 2025")

add_content_slide(prs, "Introduction", [
    "Why pathology residents should care about software: LIS limitations, workflow inefficiencies, unmet needs.",
    "Custom software can improve diagnostics, communication, and safety.",
    "Real-world impact: frozen section alerts, digital image management, urgent case triage."
])

add_content_slide(prs, "Why Build Custom Software?", [
    "Tailored to your lab's needs and workflow — especially important for subspecialty services.",
    "Enables integration with external systems (EHRs, image viewers, AI pipelines).",
    "Supports innovation and regulatory adaptation (e.g., NGS reporting, frozen sections).",
    "From Gershkovich & Sinard: Functionality gaps arise when needed tech exists but isn't in your LIS."
])

add_content_slide(prs, "Challenges to Be Aware Of", [
    "Development effort requires domain understanding.",
    "Validation and long-term support must be considered.",
    "Institutional inertia: \"we're not in the software business\" is often a myth.",
    "From Sinard & Gershkovich: Even one or two developers with domain knowledge can be effective."
])

add_content_slide(prs, "Software Lifecycle & DevSecOps in Pathology", [
    "1. Requirement Gathering – Start with specific use cases from the clinical workflow.",
    "2. Functional Specification – What should the system do? Create screen diagrams, button behaviors, exceptions.",
    "3. Technical Specification – Database, architecture, and performance planning.",
    "4. Implementation and Testing – Use version control, containers (Docker), staging environments.",
    "5. Security and Compliance – HIPAA, access logs, secure APIs.",
    "6. Deployment & Feedback Loop – Track usage and iterate."
])

add_content_slide(prs, "DevSecOps Philosophy", [
    "Secure-by-design.",
    "Rapid, adaptive iterations.",
    "Infrastructure-as-code and automation.",
    "Key point: LIS vendors often cannot move fast enough due to outdated architecture."
])

add_content_slide(prs, "Image Analysis & AI", [
    "QPath – Use Groovy scripts to quantify cell populations, analyze TILs, or export tiles for deep learning.",
    "ImageJ/Fiji – Open-source, powerful for smaller ROI and pixel-level analysis.",
    "HistomicsTK – Feature extraction and visualization for whole slide images."
])

add_content_slide(prs, "Digital Pathology Infrastructure", [
    "Digital Slide Archive (DSA) – WSI management with REST APIs.",
    "OpenSeadragon – Lightweight JS-based WSI viewer you can embed in web tools.",
    "Open Source Repositories – GitHub, SourceForge, MONAI."
])

add_content_slide(prs, "Clinical Applications", [
    "Alerting systems for STAT/RUSH cases.",
    "Specimen tracking and intraoperative diagnosis tools.",
    "NGS variant interpretation modules."
])

add_content_slide(prs, "Spotting Gaps in Workflow", [
    "What is frustrating in your daily workflow? What leads to delays, errors, or repetitive effort?",
    "From lecture material: If your LIS doesn't notify you of STAT slides without a resident, that's a gap."
])

add_content_slide(prs, "Solutions", [
    "7 Options: Ignore, patch manually, switch LIS, request vendor fix, outsource, build yourself, or integrate OSS.",
    "Residents can help design or prototype solutions with Jupyter Notebooks, low-code tools.",
    "Learn to write or read functional specs to guide developers."
])

add_content_slide(prs, "Conclusion and Call to Action", [
    "Pathology informatics is a core discipline – not optional.",
    "Custom software development is powerful, cost-effective, and sustainable.",
    "You can participate by reporting gaps, helping design use cases, or even coding solutions.",
    "Every meaningful tool starts with someone who understood the problem deeply."
])

add_content_slide(prs, "Q&A and Demo Time", [
    "Live demo of QPath on an H&E slide.",
    "Jupyter notebook using DSA API to pull case info.",
    "Discussion: What bugs you the most in your LIS today — and what would you build to fix it?"
])

# Save presentation
prs.save('../data/software_dev_clinical_use.pptx')

print("Presentation created successfully: ../data/software_dev_clinical_use.pptx")